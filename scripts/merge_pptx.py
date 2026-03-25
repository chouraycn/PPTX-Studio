"""Merge multiple PPTX files into one.

Combines slides from two or more presentations in the order specified.
Each source file's slides are appended in sequence to a fresh output file.
The first file's slide size (width × height) is used for the output.

Usage:
    python scripts/merge_pptx.py <file1.pptx> <file2.pptx> [file3.pptx ...] -o <output.pptx>

Options:
    -o, --output      Output file path (required)
    --slides          Slide range to include from each file, e.g. "1-3" or "2,4,5"
                      Apply globally to all files, or per-file (comma-separated ranges,
                      one per input file). Default: all slides.
    --order           Custom slide ordering across files. Each token is "<file_index>:<slide_num>",
                      where file_index is 1-based (1 = first file, 2 = second file, ...).
                      When --order is used, --slides is ignored.
    --ignore-notes    Do not copy speaker notes from source slides
    --dry-run         Print what would be merged without writing output

Examples:

    # Merge all slides from two files (default order: A1,A2,...,B1,B2,...)
    python scripts/merge_pptx.py a.pptx b.pptx -o merged.pptx

    # Take slides 1-5 from a.pptx and slides 2,3,7 from b.pptx
    python scripts/merge_pptx.py a.pptx b.pptx -o merged.pptx --slides "1-5" "2,3,7"

    # Interleave: A1, B1, A2, B2, A3
    python scripts/merge_pptx.py a.pptx b.pptx -o merged.pptx --order 1:1 2:1 1:2 2:2 1:3

    # Custom order: B2, A1, B1, A3
    python scripts/merge_pptx.py a.pptx b.pptx -o merged.pptx --order 2:2 1:1 2:1 1:3

    # Merge three files, skip notes
    python scripts/merge_pptx.py a.pptx b.pptx c.pptx -o merged.pptx --ignore-notes
"""

from __future__ import annotations

import argparse
import copy
import re
import sys
from pathlib import Path
from typing import Generator, Optional

# python-pptx
try:
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError:
    print("Error: python-pptx is required. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

# auto_resize module
try:
    from auto_resize import detect_slide_size, calculate_resize_strategy, resize_slide_xml, SlideSize
except ImportError:
    print("Error: auto_resize module not found. Ensure auto_resize.py is in the same directory", file=sys.stderr)
    sys.exit(1)


# ---------------------------------------------------------------------------
# Slide range parsing
# ---------------------------------------------------------------------------

def parse_order_spec(
    tokens: list[str], num_files: int, totals: list[int]
) -> list[tuple[int, int]]:
    """Parse --order tokens into a list of (0-based file_index, 0-based slide_index).

    Each token is "<file_num>:<slide_num>" where both numbers are 1-based.
    Example: ["1:1", "2:1", "1:2"] → [(0,0), (1,0), (0,1)]
    """
    result: list[tuple[int, int]] = []
    for token in tokens:
        token = token.strip()
        m = re.fullmatch(r"(\d+):(\d+)", token)
        if not m:
            raise ValueError(
                f"Invalid --order token '{token}'. Expected format: <file_num>:<slide_num>"
            )
        file_num = int(m.group(1))
        slide_num = int(m.group(2))

        if file_num < 1 or file_num > num_files:
            raise ValueError(
                f"File index {file_num} in --order token '{token}' is out of range "
                f"(1–{num_files})"
            )
        file_idx = file_num - 1
        total = totals[file_idx]
        if slide_num < 1 or slide_num > total:
            raise ValueError(
                f"Slide number {slide_num} in --order token '{token}' is out of range "
                f"for file {file_num} (1–{total})"
            )
        result.append((file_idx, slide_num - 1))
    return result


def parse_slide_range(spec: str, total: int) -> list[int]:
    """Parse a slide range spec like "1-3" or "1,3,5" into 0-based indices.

    Slide numbers are 1-based in the spec, returned as 0-based indices.
    """
    indices: list[int] = []
    for part in spec.split(","):
        part = part.strip()
        m = re.fullmatch(r"(\d+)\s*-\s*(\d+)", part)
        if m:
            start, end = int(m.group(1)), int(m.group(2))
        elif re.fullmatch(r"\d+", part):
            start = end = int(part)
        else:
            raise ValueError(f"Invalid slide range token: '{part}'")

        if start < 1 or end > total or start > end:
            raise ValueError(
                f"Slide range {start}-{end} out of bounds for a {total}-slide file"
            )
        indices.extend(range(start - 1, end))  # convert to 0-based
    return indices


# ---------------------------------------------------------------------------
# Core merge logic
# ---------------------------------------------------------------------------

def _iter_slides(prs: Presentation, indices: list[int]) -> Generator:
    """Yield (slide_element, notes_element_or_None) for each requested slide."""
    slide_list = prs.slides
    for idx in indices:
        slide = slide_list[idx]
        notes_el = None
        if slide.has_notes_slide:
            notes_el = slide.notes_slide._element
        yield slide._element, notes_el


def _append_slide(dest_prs: Presentation, slide_el: etree._Element,
                  notes_el, ignore_notes: bool) -> None:
    """Deep-copy a slide XML element into dest_prs."""
    # Blank slide template approach: add a blank slide then replace its XML
    slide_layout = dest_prs.slide_layouts[0]  # title slide layout (placeholder only)
    new_slide = dest_prs.slides.add_slide(slide_layout)

    # Replace the slide XML tree with a deep copy of the source slide
    new_slide_el = new_slide._element
    new_slide_el.getparent().replace(new_slide_el, copy.deepcopy(slide_el))

    # Re-fetch the freshly placed slide element
    placed_el = dest_prs.slides[-1]._element

    # Attach notes if requested
    if not ignore_notes and notes_el is not None:
        notes_copy = copy.deepcopy(notes_el)
        ns_slide = dest_prs.slides[-1].notes_slide
        ns_el = ns_slide._element
        ns_el.getparent().replace(ns_el, notes_copy)


def merge(
    input_files: list[str],
    output_file: str,
    slides_specs: list[str] | None = None,
    order_specs: list[str] | None = None,
    ignore_notes: bool = False,
    dry_run: bool = False,
    resize_strategy: str = "smart",
    target_size: Optional[str] = None,
    resize_warnings: bool = False,
) -> str:
    """Merge input_files into output_file.

    When order_specs is provided (list of "<file_num>:<slide_num>" tokens),
    slides are emitted in that exact custom order; slides_specs is ignored.

    Args:
        resize_strategy: Resize strategy (smart/scale/stretch/crop)
        target_size: Target slide size (16:9/4:3/16:10/auto)
        resize_warnings: Show resize warnings

    Returns a human-readable summary string.
    """
    paths = [Path(f) for f in input_files]

    # Validate inputs
    for p in paths:
        if not p.exists():
            return f"Error: {p} does not exist"
        if p.suffix.lower() != ".pptx":
            return f"Error: {p} is not a .pptx file"

    # Load presentations
    presentations = [Presentation(str(p)) for p in paths]
    totals = [len(prs.slides) for prs in presentations]

    # -----------------------------------------------------------------------
    # Mode A: --order (custom cross-file ordering)
    # -----------------------------------------------------------------------
    if order_specs:
        try:
            ordered_pairs = parse_order_spec(order_specs, len(paths), totals)
        except ValueError as e:
            return f"Error: {e}"

        # Build summary
        lines = []
        for file_idx, slide_idx in ordered_pairs:
            lines.append(
                f"  {paths[file_idx].name} slide {slide_idx + 1}"
            )
        summary = (
            f"Merging {len(paths)} files → {output_file} (custom order)\n"
            + "\n".join(lines)
            + f"\nTotal: {len(ordered_pairs)} slides"
        )

        if dry_run:
            return f"[Dry run] {summary}"

        dest_prs = Presentation()
        dest_prs.slide_width = presentations[0].slide_width
        dest_prs.slide_height = presentations[0].slide_height

        for file_idx, slide_idx in ordered_pairs:
            prs = presentations[file_idx]
            slide = prs.slides[slide_idx]
            notes_el = slide.notes_slide._element if slide.has_notes_slide else None
            _append_slide(dest_prs, slide._element, notes_el, ignore_notes)

        out_path = Path(output_file)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        dest_prs.save(str(out_path))
        return f"Done. {summary}"

    # -----------------------------------------------------------------------
    # Mode B: sequential merge (--slides optional range filter)
    # -----------------------------------------------------------------------
    resolved: list[list[int]] = []
    for i, prs in enumerate(presentations):
        total = totals[i]
        if slides_specs and i < len(slides_specs):
            spec = slides_specs[i]
        elif slides_specs and len(slides_specs) == 1:
            spec = slides_specs[0]  # single global spec
        else:
            spec = None

        if spec:
            try:
                indices = parse_slide_range(spec, total)
            except ValueError as e:
                return f"Error in slide range for {paths[i].name}: {e}"
        else:
            indices = list(range(total))

        resolved.append(indices)

    # Summary
    lines = []
    total_slides = 0
    resize_lines = []

    # Detect sizes and check if resize needed
    source_sizes = [detect_slide_size(str(p)) for p in paths]
    first_size = source_sizes[0]

    # Determine target size
    if target_size and target_size != "auto":
        from auto_resize import parse_size_spec
        target_width, target_height = parse_size_spec(target_size)
        target_size_obj = SlideSize(target_width, target_height, target_size)
    else:
        target_size_obj = first_size

    # Check if resize is needed
    needs_resize = any(
        s.aspect_ratio != target_size_obj.aspect_ratio for s in source_sizes
    )

    for p, indices, size in zip(paths, resolved, source_sizes):
        count = len(indices)
        total_slides += count
        slide_nums = [i + 1 for i in indices]
        
        line = f"  {p.name}: slides {slide_nums} ({count} slides)"
        
        if needs_resize and size.aspect_ratio != target_size_obj.aspect_ratio:
            resize_info = calculate_resize_strategy(size, target_size_obj, resize_strategy)
            line += f" [{size.aspect_ratio} → {target_size_obj.aspect_ratio}]"
            
            if resize_warnings:
                for w in resize_info.warnings:
                    resize_lines.append(f"    ⚠️  {p.name}: {w}")
        
        lines.append(line)

    summary = (
        f"Merging {len(paths)} files → {output_file}\n"
        + "\n".join(lines)
        + f"\nTotal: {total_slides} slides"
    )

    if needs_resize:
        summary += f"\nResize strategy: {resize_strategy}"
        if resize_warnings and resize_lines:
            summary += "\nResize warnings:\n" + "\n".join(resize_lines)

    if dry_run:
        return f"[Dry run] {summary}"

    # Build output presentation
    dest_prs = Presentation()
    dest_prs.slide_width = target_size_obj.width
    dest_prs.slide_height = target_size_obj.height

    # Append slides with resize if needed
    for prs, indices, source_size in zip(presentations, resolved, source_sizes):
        for slide_el, notes_el in _iter_slides(prs, indices):
            # Check if resize is needed for this source
            if source_size.aspect_ratio != target_size_obj.aspect_ratio:
                # Get slide XML as string
                slide_xml_str = etree.tostring(slide_el, encoding="unicode")
                
                # Resize the slide
                resized_xml_str = resize_slide_xml(
                    slide_xml_str,
                    source_size,
                    target_size_obj,
                    strategy=resize_strategy,
                    verbose=False
                )
                
                # Parse back to element
                resized_el = etree.fromstring(resized_xml_str)
                _append_slide(dest_prs, resized_el, notes_el, ignore_notes)
            else:
                # No resize needed, append directly
                _append_slide(dest_prs, slide_el, notes_el, ignore_notes)

    # Write output
    out_path = Path(output_file)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    dest_prs.save(str(out_path))

    return f"Done. {summary}"


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Merge multiple PPTX files into one",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    parser.add_argument(
        "files",
        nargs="+",
        metavar="FILE",
        help="Two or more .pptx files to merge (in order)",
    )
    parser.add_argument(
        "-o", "--output",
        required=True,
        metavar="OUTPUT",
        help="Output .pptx file path",
    )
    parser.add_argument(
        "--slides",
        nargs="+",
        metavar="RANGE",
        default=None,
        help=(
            'Slide range(s) to include. One value applies to all files; '
            'multiple values apply per-file in order. '
            'Format: "1-3" (range) or "1,3,5" (list). 1-based numbering. '
            'Ignored when --order is used.'
        ),
    )
    parser.add_argument(
        "--order",
        nargs="+",
        metavar="FILE:SLIDE",
        default=None,
        help=(
            'Custom slide ordering. Each token is "<file_num>:<slide_num>" (both 1-based). '
            'Defines the exact sequence of slides in the output, enabling interleaving '
            'or arbitrary reordering across source files. '
            'Example: --order 1:1 2:1 1:2 2:2  →  A1, B1, A2, B2'
        ),
    )
    parser.add_argument(
        "--ignore-notes",
        action="store_true",
        help="Do not copy speaker notes from source slides",
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Print merge plan without writing output",
    )
    parser.add_argument(
        "--resize-strategy",
        choices=["smart", "scale", "stretch", "crop"],
        default="smart",
        help=(
            "Resize strategy when merging different slide sizes. "
            "smart: AI layout optimization (default), "
            "scale: maintain aspect ratio and center, "
            "stretch: stretch to fill (may distort), "
            "crop: crop to fill (may lose content)"
        ),
    )
    parser.add_argument(
        "--target-size",
        choices=["16:9", "4:3", "16:10", "auto"],
        default="auto",
        help=(
            "Target slide size for merged output. "
            "auto: use first file's size (default), "
            "16:9/4:3/16:10: force specific size"
        ),
    )
    parser.add_argument(
        "--resize-warnings",
        action="store_true",
        help="Show warnings about resize operations",
    )

    args = parser.parse_args()

    if len(args.files) < 2:
        print("Error: at least two input files are required", file=sys.stderr)
        sys.exit(1)

    result = merge(
        input_files=args.files,
        output_file=args.output,
        slides_specs=args.slides,
        order_specs=args.order,
        ignore_notes=args.ignore_notes,
        dry_run=args.dry_run,
        resize_strategy=args.resize_strategy,
        target_size=args.target_size,
        resize_warnings=args.resize_warnings,
    )

    print(result)

    if result.startswith("Error"):
        sys.exit(1)
