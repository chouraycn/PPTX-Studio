#!/usr/bin/env python3
"""
generate_notes.py — 自动为 PPTX 每页幻灯片生成演讲者备注

功能：
  1. 读取每页幻灯片的标题、正文内容
  2. 调用 AI（本地 LLM 或 OpenAI API）生成对应备注
  3. 将备注写入每页幻灯片的 notes 字段
  4. 支持"仅追加"模式（已有备注的页不覆盖）

用法：
  python scripts/generate_notes.py input.pptx output.pptx [选项]

  --mode         备注风格：speaker（演讲提示）| coach（演讲教练）| summary（摘要）
                 默认: speaker
  --language     备注语言：zh（中文）| en（英文）| auto（随内容自动检测）
                 默认: auto
  --no-overwrite   跳过已有备注的幻灯片（默认覆盖）
  --append-summary 保留原有备注，在其后追加 AI 总结分区（推荐用于模板套用后增强备注）
  --dry-run        只打印生成的备注，不写入文件
  --backend      生成后端：openai | ollama | simple
                 默认: simple（无需 API，基于规则生成）
  --model        模型名称，仅 openai/ollama 后端有效
  --api-key      OpenAI API Key（也可通过 OPENAI_API_KEY 环境变量设置）
  --base-url     API Base URL（用于 ollama 或自定义端点）

示例：
  # 最简单用法（规则生成，无需 API）
  python scripts/generate_notes.py deck.pptx deck_with_notes.pptx

  # 使用 OpenAI 生成
  python scripts/generate_notes.py deck.pptx out.pptx --backend openai --api-key sk-xxx

  # 使用本地 Ollama
  python scripts/generate_notes.py deck.pptx out.pptx --backend ollama --model llama3

  # 仅打印，不写入
  python scripts/generate_notes.py deck.pptx out.pptx --dry-run

  # 跳过已有备注的页
  python scripts/generate_notes.py deck.pptx out.pptx --no-overwrite

  # 保留原有备注，追加 AI 总结（模板套用后增强备注的推荐用法）
  python scripts/generate_notes.py output.pptx output_with_notes.pptx --append-summary

  # 保留原有备注 + 追加 AI 总结，使用 OpenAI 后端
  python scripts/generate_notes.py output.pptx output_with_notes.pptx --append-summary --backend openai --api-key sk-xxx
"""

import argparse
import json
import os
import re
import sys
import textwrap
from pathlib import Path
from typing import Optional

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("错误：需要安装 python-pptx")
    print("  pip install python-pptx")
    sys.exit(1)

# ──────────────────────────────────────────────
# 辅助：提取幻灯片内容
# ──────────────────────────────────────────────

def extract_slide_content(slide) -> dict:
    """提取单张幻灯片的文本内容。"""
    title = ""
    body_lines = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        # 判断是否是标题占位符
        is_title = (
            hasattr(shape, "placeholder_format")
            and shape.placeholder_format is not None
            and shape.placeholder_format.idx in (0, 1)  # 0=title, 1=center title
        )
        if is_title:
            title = text
        else:
            # 拆分成行
            for line in text.splitlines():
                line = line.strip()
                if line:
                    body_lines.append(line)

    # 检测是否有图表/图片
    has_chart = any(shape.has_chart for shape in slide.shapes if hasattr(shape, "has_chart"))
    has_image = any(
        shape.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
        for shape in slide.shapes
    )

    return {
        "title": title,
        "body": body_lines,
        "has_chart": has_chart,
        "has_image": has_image,
    }


def get_existing_notes(slide) -> str:
    """获取幻灯片已有的备注文字。"""
    if slide.has_notes_slide:
        tf = slide.notes_slide.notes_text_frame
        return tf.text.strip()
    return ""


# ──────────────────────────────────────────────
# 语言检测（简单规则）
# ──────────────────────────────────────────────

def detect_language(text: str) -> str:
    """简单判断文本主要语言：zh 或 en。"""
    chinese_chars = sum(1 for c in text if "\u4e00" <= c <= "\u9fff")
    if chinese_chars > len(text) * 0.1:
        return "zh"
    return "en"


# ──────────────────────────────────────────────
# 备注生成 — simple 后端（规则驱动，无 API）
# ──────────────────────────────────────────────

TRANSITION_ZH = [
    "接下来我们来看……",
    "让我们深入了解……",
    "这里有一个关键点……",
    "请注意……",
    "我想强调的是……",
    "正如大家看到的……",
    "这背后的逻辑是……",
]

TRANSITION_EN = [
    "Let's take a closer look at...",
    "Here I want to highlight...",
    "Notice that...",
    "The key takeaway here is...",
    "As you can see...",
    "Let me walk you through...",
    "The reasoning behind this is...",
]

CLOSING_ZH = [
    "如果大家有问题，我们稍后讨论。",
    "这一页内容比较多，欢迎随时提问。",
    "记住这个核心信息：{title}。",
]

CLOSING_EN = [
    "Feel free to ask questions after this section.",
    "This is a key slide — take a moment to absorb it.",
    "Remember the core message: {title}.",
]

def _pick(lst, idx):
    return lst[idx % len(lst)]

def generate_notes_simple(slide_content: dict, slide_idx: int, mode: str, lang: str) -> str:
    """
    规则驱动生成备注，无需任何 API。
    适合快速使用或离线场景。
    """
    title = slide_content["title"]
    body = slide_content["body"]
    has_chart = slide_content["has_chart"]
    has_image = slide_content["has_image"]

    if lang == "auto":
        combined = title + " ".join(body)
        lang = detect_language(combined) if combined.strip() else "zh"

    if lang == "zh":
        transitions = TRANSITION_ZH
        closing_templates = CLOSING_ZH
        chart_note = "这里的图表展示了核心数据趋势，可以结合实际业务数字进行说明。"
        image_note = "配图帮助听众直观理解内容，可根据现场情况选择是否展开讲解。"
        no_content_note = "（此页为过渡/章节页，可简短过渡，引出下一部分。）"
    else:
        transitions = TRANSITION_EN
        closing_templates = CLOSING_EN
        chart_note = "This chart illustrates the core data trend. Feel free to highlight specific numbers relevant to your audience."
        image_note = "The visual supports the narrative — expand on it as needed based on audience engagement."
        no_content_note = "(This is a transition or section header slide. Keep it brief and lead into the next part.)"

    parts = []

    # ── 模式：speaker（演讲提示）──
    if mode == "speaker":
        # 开场过渡句
        if title:
            trans = _pick(transitions, slide_idx)
            if lang == "zh":
                parts.append(f"{trans}「{title}」")
            else:
                parts.append(f"{trans} \"{title}\"")

        # 逐条要点提示
        if body:
            if lang == "zh":
                parts.append("\n关键要点提示：")
            else:
                parts.append("\nKey talking points:")
            for i, line in enumerate(body[:5], 1):
                # 截短超长行
                short = (line[:60] + "…") if len(line) > 60 else line
                if lang == "zh":
                    parts.append(f"  {i}. {short} — 展开说明背景或数据支撑。")
                else:
                    parts.append(f"  {i}. {short} — Elaborate with context or data.")

        # 图表/图片提示
        if has_chart:
            parts.append(f"\n{chart_note}")
        if has_image:
            parts.append(f"\n{image_note}")

        # 无内容页
        if not title and not body:
            parts.append(no_content_note)

        # 收尾
        if title:
            closing = _pick(closing_templates, slide_idx).replace("{title}", title)
            parts.append(f"\n{closing}")

    # ── 模式：coach（演讲教练）──
    elif mode == "coach":
        if lang == "zh":
            parts.append("【演讲教练提示】")
            if title:
                parts.append(f"本页主题：{title}")
            if len(body) > 5:
                parts.append(f"⚠️ 内容偏多（{len(body)} 条）——建议口头精简，只讲核心 3 点。")
            elif len(body) == 0 and title:
                parts.append("本页文字较少，注意填充口头内容，避免空洞。")
            else:
                parts.append(f"内容量适中（{len(body)} 条），正常节奏即可。")
            parts.append("\n建议用时：约 1-2 分钟。")
            parts.append("转场建议：结束本页时给听众留 3 秒停顿，再进入下一页。")
        else:
            parts.append("[Coaching Notes]")
            if title:
                parts.append(f"Slide topic: {title}")
            if len(body) > 5:
                parts.append(f"⚠️ Too many points ({len(body)}) — verbally trim to top 3.")
            elif len(body) == 0 and title:
                parts.append("Sparse content — pad with verbal context to avoid gaps.")
            else:
                parts.append(f"Good content density ({len(body)} points) — proceed at normal pace.")
            parts.append("\nSuggested time: ~1-2 minutes.")
            parts.append("Transition tip: Pause 3 seconds after this slide before advancing.")

    # ── 模式：summary（摘要）──
    elif mode == "summary":
        if lang == "zh":
            if title:
                parts.append(f"本页核心：{title}")
            if body:
                parts.append("要点：" + "；".join(b[:40] for b in body[:4]) + "。")
            if has_chart:
                parts.append("（含图表）")
            if has_image:
                parts.append("（含配图）")
        else:
            if title:
                parts.append(f"Core: {title}")
            if body:
                parts.append("Points: " + "; ".join(b[:40] for b in body[:4]) + ".")
            if has_chart:
                parts.append("(Chart included)")
            if has_image:
                parts.append("(Image included)")

    return "\n".join(parts).strip()


# ──────────────────────────────────────────────
# 备注生成 — OpenAI 后端
# ──────────────────────────────────────────────

def generate_notes_openai(
    slide_content: dict,
    slide_idx: int,
    mode: str,
    lang: str,
    api_key: str,
    model: str = "gpt-4o-mini",
    base_url: Optional[str] = None,
) -> str:
    try:
        import openai
    except ImportError:
        print("错误：需要安装 openai 包")
        print("  pip install openai")
        sys.exit(1)

    title = slide_content["title"]
    body = slide_content["body"]

    content_desc = f"标题：{title}\n正文：\n" + "\n".join(f"- {b}" for b in body)
    if not (title or body):
        content_desc = "（无文字内容，可能是图表页或过渡页）"

    mode_instructions = {
        "speaker": "请生成实用的演讲提示，包含：开场过渡句、每个要点的口头展开建议、收尾句。",
        "coach": "请以演讲教练角度提供反馈：内容密度是否合适、建议用时、过渡技巧。",
        "summary": "请用一到两句话概括本页核心信息。",
    }

    lang_instruction = ""
    if lang == "zh":
        lang_instruction = "请用中文生成备注。"
    elif lang == "en":
        lang_instruction = "Please generate notes in English."
    else:
        lang_instruction = "请根据幻灯片内容语言自动选择备注语言。"

    prompt = f"""你是一名资深演示顾问，正在为演示者准备演讲者备注。

幻灯片 {slide_idx + 1} 内容：
{content_desc}

{mode_instructions.get(mode, mode_instructions['speaker'])}
{lang_instruction}

要求：
- 简洁实用，控制在 150 字以内（英文 100 词以内）
- 不要复述幻灯片原文，而是提供额外的口头价值
- 不要使用 Markdown 格式，输出纯文本
"""

    client_kwargs = {"api_key": api_key}
    if base_url:
        client_kwargs["base_url"] = base_url

    client = openai.OpenAI(**client_kwargs)
    response = client.chat.completions.create(
        model=model,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7,
        max_tokens=300,
    )
    return response.choices[0].message.content.strip()


# ──────────────────────────────────────────────
# 备注生成 — Ollama 后端
# ──────────────────────────────────────────────

def generate_notes_ollama(
    slide_content: dict,
    slide_idx: int,
    mode: str,
    lang: str,
    model: str = "llama3",
    base_url: str = "http://localhost:11434",
) -> str:
    try:
        import urllib.request
        import json as _json
    except ImportError:
        pass

    title = slide_content["title"]
    body = slide_content["body"]
    content_desc = f"Title: {title}\nBody:\n" + "\n".join(f"- {b}" for b in body)

    prompt = f"""You are a presentation coach. Generate concise speaker notes for slide {slide_idx + 1}.

{content_desc}

Mode: {mode}
Language: {"Chinese" if lang == "zh" else "English" if lang == "en" else "auto-detect"}

Keep it under 100 words. No markdown. Plain text only."""

    payload = json.dumps({"model": model, "prompt": prompt, "stream": False}).encode()
    req = urllib.request.Request(
        f"{base_url.rstrip('/')}/api/generate",
        data=payload,
        headers={"Content-Type": "application/json"},
    )
    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            data = json.loads(resp.read())
            return data.get("response", "").strip()
    except Exception as e:
        print(f"  ⚠️  Ollama 调用失败: {e}，回退到 simple 模式")
        return generate_notes_simple(slide_content, slide_idx, mode, lang)


# ──────────────────────────────────────────────
# 写入备注到幻灯片
# ──────────────────────────────────────────────

def write_notes_to_slide(slide, notes_text: str, append_summary: bool = False, lang: str = "zh"):
    """将备注文字写入幻灯片的 notes placeholder。

    Args:
        slide: python-pptx Slide 对象
        notes_text: 要写入的备注文本（新生成的 AI 内容）
        append_summary: 若为 True，保留原有备注，在其后追加分隔线 + AI 总结；
                        若为 False，覆盖原有备注（原行为）。
        lang: 分隔线语言，"zh" 中文 / "en" 英文
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    if not slide.has_notes_slide:
        _ = slide.notes_slide

    tf = slide.notes_slide.notes_text_frame
    txBody = tf._txBody

    if append_summary:
        # ── 追加模式：保留原有备注，添加分隔线 + AI 总结 ──
        # 分隔符文本
        separator = "────────────────────────" if lang == "en" else "────────────────────────"
        if lang == "zh":
            section_header = "【AI 摘要】"
        else:
            section_header = "[AI Summary]"

        # 追加分隔线段落
        for line_text in ["", separator, section_header, ""]:
            new_p = etree.SubElement(txBody, qn("a:p"))
            if line_text:
                new_r = etree.SubElement(new_p, qn("a:r"))
                new_t = etree.SubElement(new_r, qn("a:t"))
                new_t.text = line_text

        # 追加 AI 总结内容（逐行）
        for line in notes_text.split("\n"):
            new_p = etree.SubElement(txBody, qn("a:p"))
            if line.strip():
                new_r = etree.SubElement(new_p, qn("a:r"))
                new_t = etree.SubElement(new_r, qn("a:t"))
                new_t.text = line
    else:
        # ── 覆盖模式（原行为）：彻底删除所有旧 <a:p> 节点后重写 ──
        for old_p in txBody.findall(qn("a:p")):
            txBody.remove(old_p)

        for line in notes_text.split("\n"):
            new_p = etree.SubElement(txBody, qn("a:p"))
            new_r = etree.SubElement(new_p, qn("a:r"))
            new_t = etree.SubElement(new_r, qn("a:t"))
            new_t.text = line


# ──────────────────────────────────────────────
# 主流程
# ──────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="为 PPTX 每页幻灯片自动生成演讲者备注",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    parser.add_argument("input", help="输入 PPTX 文件路径")
    parser.add_argument("output", help="输出 PPTX 文件路径")
    parser.add_argument(
        "--mode",
        choices=["speaker", "coach", "summary"],
        default="speaker",
        help="备注风格（默认: speaker）",
    )
    parser.add_argument(
        "--language",
        choices=["zh", "en", "auto"],
        default="auto",
        help="备注语言（默认: auto）",
    )
    parser.add_argument(
        "--no-overwrite",
        action="store_true",
        help="跳过已有备注的幻灯片",
    )
    parser.add_argument(
        "--append-summary",
        action="store_true",
        help="保留原有备注，在其后追加分隔线 + AI 总结（模板套用后增强备注的推荐用法）",
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="只打印备注，不写入文件",
    )
    parser.add_argument(
        "--backend",
        choices=["simple", "openai", "ollama"],
        default="simple",
        help="生成后端（默认: simple）",
    )
    parser.add_argument("--model", default=None, help="模型名称（openai/ollama 有效）")
    parser.add_argument("--api-key", default=None, help="OpenAI API Key")
    parser.add_argument("--base-url", default=None, help="API Base URL")

    args = parser.parse_args()

    # 验证输入文件
    input_path = Path(args.input)
    if not input_path.exists():
        print(f"错误：文件不存在 — {input_path}")
        sys.exit(1)

    # OpenAI API Key
    api_key = args.api_key or os.environ.get("OPENAI_API_KEY", "")
    if args.backend == "openai" and not api_key:
        print("错误：openai 后端需要提供 --api-key 或设置 OPENAI_API_KEY 环境变量")
        sys.exit(1)

    # 默认模型
    default_models = {"openai": "gpt-4o-mini", "ollama": "llama3"}
    model = args.model or default_models.get(args.backend, "")

    print(f"📂 加载: {input_path}")
    prs = Presentation(str(input_path))
    total = len(prs.slides)
    mode_desc = "追加AI摘要" if args.append_summary else args.mode
    print(f"📊 共 {total} 张幻灯片 | 模式: {mode_desc} | 语言: {args.language} | 后端: {args.backend}")
    print()

    generated_count = 0
    skipped_count = 0

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        content = extract_slide_content(slide)
        existing = get_existing_notes(slide)

        # 跳过已有备注的页（append-summary 模式下不跳过，因为追加不会破坏原有备注）
        if args.no_overwrite and existing and not args.append_summary:
            print(f"  [slide {slide_num:2d}] ⏭  跳过（已有备注）")
            skipped_count += 1
            continue

        # append-summary 模式下，若该页已追加过 AI 摘要，跳过（避免重复追加）
        if args.append_summary and existing and ("【AI 摘要】" in existing or "[AI Summary]" in existing):
            print(f"  [slide {slide_num:2d}] ⏭  跳过（已含 AI 摘要）")
            skipped_count += 1
            continue

        # 生成备注
        print(f"  [slide {slide_num:2d}] ✍  生成中... 标题: {content['title'][:30] or '(无标题)'}")

        # append-summary 模式下固定使用 summary 风格，内容更精简
        note_mode = "summary" if args.append_summary else args.mode

        try:
            if args.backend == "openai":
                notes = generate_notes_openai(
                    content, i, note_mode, args.language, api_key, model, args.base_url
                )
            elif args.backend == "ollama":
                notes = generate_notes_ollama(
                    content, i, note_mode, args.language, model,
                    args.base_url or "http://localhost:11434"
                )
            else:
                notes = generate_notes_simple(content, i, note_mode, args.language)
        except Exception as e:
            print(f"    ⚠️  生成失败: {e}，使用 simple 后端兜底")
            notes = generate_notes_simple(content, i, note_mode, args.language)

        # 检测语言（用于分隔符语言选择）
        combined_text = content["title"] + " ".join(content["body"])
        detected_lang = detect_language(combined_text) if combined_text.strip() else "zh"
        if args.language != "auto":
            detected_lang = args.language

        if args.dry_run:
            if args.append_summary and existing:
                print(f"    ── 原有备注 ──")
                for line in existing.splitlines()[:3]:
                    print(f"    {line}")
                print(f"    ...")
            print(f"    ── AI 摘要追加 ──" if args.append_summary else f"    ── 备注预览 ──")
            for line in notes.splitlines():
                print(f"    {line}")
            print()
        else:
            write_notes_to_slide(
                slide, notes,
                append_summary=args.append_summary,
                lang=detected_lang,
            )

        generated_count += 1

    # 保存
    if not args.dry_run:
        output_path = Path(args.output)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        print()
        print(f"✅ 完成！已生成 {generated_count} 页备注，跳过 {skipped_count} 页")
        print(f"💾 已保存: {output_path}")
    else:
        print(f"🔍 dry-run 完成：已预览 {generated_count} 页（未写入文件）")


if __name__ == "__main__":
    main()
