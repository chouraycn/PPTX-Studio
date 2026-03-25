"""Automated QA checker for PPTX files.

Runs a battery of structural and visual quality checks on a presentation
without requiring LibreOffice or any external tools. Pure python-pptx.

Checks performed:
  - Text overflow: text frame content exceeds its bounding box
  - Contrast: text color vs background color (WCAG AA: 4.5:1 for body, 3:1 for large)
  - Empty slides: slides with no visible text or shapes
  - Placeholder leakage: leftover "Click to edit" / "lorem ipsum" text
  - Bullet overload: slides with more than 6 bullet points
  - Font size violations: body text below 12pt or title below 20pt
  - Off-slide elements: shapes positioned outside the slide boundaries
  - Duplicate slides: consecutive slides with identical titles
  - Missing titles: content slides with no detectable title
  - Layout monotony: 3+ consecutive slides using the same layout

Usage:
    python scripts/qa_check.py presentation.pptx
    python scripts/qa_check.py presentation.pptx --output report.json
    python scripts/qa_check.py presentation.pptx --strict        # exit 1 if issues found
    python scripts/qa_check.py presentation.pptx --only contrast,overflow
    python scripts/qa_check.py presentation.pptx --min-severity warning
"""

import argparse
import json
import re
import sys
import zipfile
from pathlib import Path
from typing import Optional, List, Dict, Any, Tuple

try:
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

SEVERITY_ERROR   = "error"    # Must fix — likely corrupted output or unreadable content
SEVERITY_WARNING = "warning"  # Should fix — degrades quality
SEVERITY_INFO    = "info"     # Nice to fix — minor polish

PLACEHOLDER_PATTERNS = [
    r"click to edit",
    r"lorem ipsum",
    r"your title here",
    r"add title",
    r"add text",
    r"text placeholder",
    r"title \d+",
    r"subtitle",
    r"xxxx",
    r"this (page|slide) layout",
    r"insert (title|text|content)",
    r"placeholder",
]

ALL_CHECKS = [
    "overflow",
    "contrast",
    "empty",
    "placeholder",
    "bullets",
    "fontsize",
    "offslide",
    "duplicates",
    "titles",
    "monotony",
]

# ---------------------------------------------------------------------------
# Color utilities
# ---------------------------------------------------------------------------

def _hex_to_rgb(hex_str: str) -> Tuple[int, int, int]:
    """Convert 6-char hex string to (r, g, b) tuple."""
    hex_str = hex_str.lstrip("#").upper()
    if len(hex_str) != 6:
        return (128, 128, 128)  # fallback gray
    try:
        r = int(hex_str[0:2], 16)
        g = int(hex_str[2:4], 16)
        b = int(hex_str[4:6], 16)
        return (r, g, b)
    except ValueError:
        return (128, 128, 128)


def _relative_luminance(r: int, g: int, b: int) -> float:
    """WCAG relative luminance (0=black, 1=white)."""
    def channel(c: int) -> float:
        s = c / 255.0
        return s / 12.92 if s <= 0.03928 else ((s + 0.055) / 1.055) ** 2.4
    return 0.2126 * channel(r) + 0.7152 * channel(g) + 0.0722 * channel(b)


def _contrast_ratio(fg: Tuple[int, int, int], bg: Tuple[int, int, int]) -> float:
    """WCAG contrast ratio between two RGB tuples."""
    l1 = _relative_luminance(*fg)
    l2 = _relative_luminance(*bg)
    lighter = max(l1, l2)
    darker  = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def _get_shape_background_color(shape) -> Optional[Tuple[int, int, int]]:
    """Best-effort extraction of a shape's fill color."""
    try:
        fill = shape.fill
        if fill.type is not None:
            from pptx.enum.dml import MSO_THEME_COLOR
            from pptx.dml.color import RGBColor as _RGB
            rgb = fill.fore_color.rgb
            if rgb:
                return (rgb.red, rgb.green, rgb.blue)
    except Exception:
        pass
    return None


def _get_slide_background_color(slide) -> Tuple[int, int, int]:
    """Extract slide background color. Falls back to white."""
    try:
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            rgb = fill.fore_color.rgb
            if rgb:
                return (rgb.red, rgb.green, rgb.blue)
    except Exception:
        pass
    return (255, 255, 255)  # default white


def _get_run_color(run) -> Optional[Tuple[int, int, int]]:
    """Extract font color from a text run."""
    try:
        rgb = run.font.color.rgb
        if rgb:
            return (rgb.red, rgb.green, rgb.blue)
    except Exception:
        pass
    return None

# ---------------------------------------------------------------------------
# Issue builder
# ---------------------------------------------------------------------------

def _issue(check: str, severity: str, slide_num: int, message: str,
           detail: Optional[str] = None) -> Dict[str, Any]:
    issue = {
        "check":    check,
        "severity": severity,
        "slide":    slide_num,
        "message":  message,
    }
    if detail:
        issue["detail"] = detail
    return issue

# ---------------------------------------------------------------------------
# Individual checks
# ---------------------------------------------------------------------------

def check_overflow(slide, slide_num: int) -> List[Dict]:
    """Detect text frames where content likely overflows the bounding box."""
    issues = []
    slide_width  = slide.shapes.width  if hasattr(slide.shapes, "width")  else None
    slide_height = slide.shapes.height if hasattr(slide.shapes, "height") else None

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        if not tf.text.strip():
            continue

        # Count paragraphs and estimate total text height
        total_lines = 0
        for para in tf.paragraphs:
            text = para.text.strip()
            if not text:
                total_lines += 1  # blank line
                continue
            # Estimate wrapping: rough chars-per-line based on shape width
            shape_width_pt = shape.width / 12700  # EMU to points (1pt = 12700 EMU)
            avg_char_width_pt = 7  # rough estimate for 14pt font
            chars_per_line = max(1, int(shape_width_pt / avg_char_width_pt))
            lines = max(1, (len(text) + chars_per_line - 1) // chars_per_line)
            total_lines += lines

        # Estimate required height
        avg_line_height_pt = 18  # line height ~18pt
        required_height_pt = total_lines * avg_line_height_pt
        shape_height_pt = shape.height / 12700

        if required_height_pt > shape_height_pt * 1.25:  # 25% tolerance
            issues.append(_issue(
                "overflow", SEVERITY_WARNING, slide_num,
                f"Text may overflow in shape '{shape.name}'",
                f"Estimated {required_height_pt:.0f}pt needed, box is {shape_height_pt:.0f}pt tall"
            ))

    return issues


def check_contrast(slide, slide_num: int) -> List[Dict]:
    """Check text/background contrast ratios."""
    issues = []
    slide_bg = _get_slide_background_color(slide)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        shape_bg = _get_shape_background_color(shape) or slide_bg

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                text = run.text.strip()
                if not text:
                    continue

                fg = _get_run_color(run)
                if fg is None:
                    # No explicit color set — assume inherits (black on white = fine)
                    continue

                ratio = _contrast_ratio(fg, shape_bg)
                font_size_pt = (run.font.size or Pt(14)) / 12700.0  # EMU to pt

                # WCAG AA: 4.5:1 for normal text, 3:1 for large text (>18pt or >14pt bold)
                is_large = (font_size_pt >= 18) or (font_size_pt >= 14 and run.font.bold)
                threshold = 3.0 if is_large else 4.5

                if ratio < threshold:
                    fg_hex = "#{:02X}{:02X}{:02X}".format(*fg)
                    bg_hex = "#{:02X}{:02X}{:02X}".format(*shape_bg)
                    issues.append(_issue(
                        "contrast", SEVERITY_WARNING, slide_num,
                        f"Low contrast in shape '{shape.name}': ratio {ratio:.1f}:1 (need {threshold}:1)",
                        f"Text color {fg_hex} on background {bg_hex} — snippet: \"{text[:40]}\""
                    ))

    return issues


def check_empty(slide, slide_num: int) -> List[Dict]:
    """Detect slides with no visible text or meaningful shapes."""
    issues = []
    has_text = False
    has_image = False

    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            has_text = True
            break
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            has_image = True

    if not has_text and not has_image:
        issues.append(_issue(
            "empty", SEVERITY_ERROR, slide_num,
            "Slide appears to be empty (no text or images)"
        ))

    return issues


def check_placeholder_leakage(slide, slide_num: int) -> List[Dict]:
    """Detect leftover template placeholder text."""
    issues = []
    pattern = re.compile("|".join(PLACEHOLDER_PATTERNS), re.IGNORECASE)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        if pattern.search(text):
            issues.append(_issue(
                "placeholder", SEVERITY_ERROR, slide_num,
                f"Leftover placeholder text in shape '{shape.name}'",
                f"Text: \"{text[:80]}\""
            ))

    return issues


def check_bullet_overload(slide, slide_num: int) -> List[Dict]:
    """Flag slides with too many bullet points."""
    issues = []
    MAX_BULLETS = 6

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Count non-empty paragraphs that look like bullet points
        bullet_count = sum(
            1 for para in shape.text_frame.paragraphs
            if para.text.strip() and para.level > 0
        )
        # Also count top-level body paras (level 0 but in a body placeholder)
        if hasattr(shape, "placeholder_format") and shape.placeholder_format:
            ph_idx = shape.placeholder_format.idx
            if ph_idx not in (0, 1):  # not title/subtitle
                bullet_count = sum(
                    1 for para in shape.text_frame.paragraphs
                    if para.text.strip()
                )

        if bullet_count > MAX_BULLETS:
            issues.append(_issue(
                "bullets", SEVERITY_WARNING, slide_num,
                f"Too many bullet points in shape '{shape.name}': {bullet_count} (max {MAX_BULLETS})",
                "Consider splitting into two slides or using a visual layout"
            ))

    return issues


def check_font_size(slide, slide_num: int) -> List[Dict]:
    """Detect font sizes that are too small for readability."""
    issues = []
    MIN_BODY_PT  = 12
    MIN_TITLE_PT = 20

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        is_title = False
        if hasattr(shape, "placeholder_format") and shape.placeholder_format:
            if shape.placeholder_format.idx in (0, 1):
                is_title = True

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                if run.font.size is None:
                    continue
                size_pt = run.font.size / 12700.0
                min_pt = MIN_TITLE_PT if is_title else MIN_BODY_PT

                if size_pt < min_pt:
                    issues.append(_issue(
                        "fontsize", SEVERITY_WARNING, slide_num,
                        f"Font too small in shape '{shape.name}': {size_pt:.0f}pt (min {min_pt}pt)",
                        f"Text: \"{run.text[:40]}\""
                    ))

    return issues


def check_offslide_elements(slide, slide_num: int, slide_width: int, slide_height: int) -> List[Dict]:
    """Detect shapes positioned outside the slide boundary."""
    issues = []

    for shape in slide.shapes:
        left   = shape.left   if shape.left   is not None else 0
        top    = shape.top    if shape.top    is not None else 0
        right  = left + (shape.width  or 0)
        bottom = top  + (shape.height or 0)

        if right < 0 or left > slide_width or bottom < 0 or top > slide_height:
            issues.append(_issue(
                "offslide", SEVERITY_WARNING, slide_num,
                f"Shape '{shape.name}' is completely outside the slide boundaries",
                f"Position: left={left//12700}pt top={top//12700}pt "
                f"right={right//12700}pt bottom={bottom//12700}pt "
                f"(slide is {slide_width//12700}pt × {slide_height//12700}pt)"
            ))
        elif left < 0 or top < 0 or right > slide_width * 1.02 or bottom > slide_height * 1.02:
            issues.append(_issue(
                "offslide", SEVERITY_INFO, slide_num,
                f"Shape '{shape.name}' partially extends outside slide boundaries",
                f"Position: left={left//12700}pt top={top//12700}pt "
                f"right={right//12700}pt bottom={bottom//12700}pt"
            ))

    return issues


def check_duplicate_slides(slides_data: List[Dict]) -> List[Dict]:
    """Detect consecutive slides with identical titles."""
    issues = []
    prev_title = None
    prev_num   = None

    for entry in slides_data:
        title = entry.get("title", "").strip().lower()
        num   = entry["num"]

        if title and title == prev_title:
            issues.append(_issue(
                "duplicates", SEVERITY_WARNING, num,
                f"Slide {num} has the same title as slide {prev_num}: \"{entry['title']}\""
            ))

        prev_title = title
        prev_num   = num

    return issues


def check_missing_titles(slides_data: List[Dict]) -> List[Dict]:
    """Detect content slides with no title."""
    issues = []
    for entry in slides_data:
        if entry.get("is_content") and not entry.get("title", "").strip():
            issues.append(_issue(
                "titles", SEVERITY_WARNING, entry["num"],
                "Content slide has no detectable title"
            ))
    return issues


def check_layout_monotony(slides_data: List[Dict]) -> List[Dict]:
    """Flag 3+ consecutive slides using the exact same layout name."""
    issues = []
    streak = 1
    prev_layout = None
    streak_start = 1

    for entry in slides_data:
        layout = entry.get("layout", "")
        num    = entry["num"]

        if layout and layout == prev_layout:
            streak += 1
            if streak == 3:
                issues.append(_issue(
                    "monotony", SEVERITY_INFO, num,
                    f"3+ consecutive slides using layout '{layout}' (started at slide {streak_start})",
                    "Vary layouts to maintain visual interest"
                ))
        else:
            streak = 1
            streak_start = num

        prev_layout = layout

    return issues

# ---------------------------------------------------------------------------
# Main runner
# ---------------------------------------------------------------------------

def run_qa(pptx_path: str, enabled_checks: Optional[List[str]] = None) -> Dict[str, Any]:
    """Run all QA checks on a PPTX file. Returns a structured report dict."""
    path = Path(pptx_path)
    if not path.exists():
        return {"error": f"File not found: {pptx_path}", "issues": [], "summary": {}}

    prs = Presentation(str(path))
    slide_width  = prs.slide_width
    slide_height = prs.slide_height

    checks = enabled_checks or ALL_CHECKS
    all_issues: List[Dict] = []
    slides_data: List[Dict] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        # Collect slide metadata for cross-slide checks
        title_shape = None
        for shape in slide.shapes:
            if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                if shape.placeholder_format.idx == 0:
                    title_shape = shape
                    break

        title_text = title_shape.text_frame.text.strip() if title_shape else ""
        has_body   = any(
            shape.has_text_frame and shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape != title_shape
        )

        slides_data.append({
            "num":        slide_num,
            "title":      title_text,
            "layout":     slide.slide_layout.name if slide.slide_layout else "",
            "is_content": bool(has_body),
        })

        # Per-slide checks
        if "overflow"    in checks: all_issues += check_overflow(slide, slide_num)
        if "contrast"    in checks: all_issues += check_contrast(slide, slide_num)
        if "empty"       in checks: all_issues += check_empty(slide, slide_num)
        if "placeholder" in checks: all_issues += check_placeholder_leakage(slide, slide_num)
        if "bullets"     in checks: all_issues += check_bullet_overload(slide, slide_num)
        if "fontsize"    in checks: all_issues += check_font_size(slide, slide_num)
        if "offslide"    in checks: all_issues += check_offslide_elements(slide, slide_num, slide_width, slide_height)

    # Cross-slide checks
    if "duplicates" in checks: all_issues += check_duplicate_slides(slides_data)
    if "titles"     in checks: all_issues += check_missing_titles(slides_data)
    if "monotony"   in checks: all_issues += check_layout_monotony(slides_data)

    # Build summary
    counts = {SEVERITY_ERROR: 0, SEVERITY_WARNING: 0, SEVERITY_INFO: 0}
    for issue in all_issues:
        counts[issue["severity"]] += 1

    return {
        "file":       str(path),
        "slides":     len(prs.slides),
        "issues":     all_issues,
        "summary": {
            "total":    len(all_issues),
            "errors":   counts[SEVERITY_ERROR],
            "warnings": counts[SEVERITY_WARNING],
            "info":     counts[SEVERITY_INFO],
            "passed":   len(all_issues) == 0,
        }
    }

# ---------------------------------------------------------------------------
# Output formatting
# ---------------------------------------------------------------------------

SEVERITY_ICON = {
    SEVERITY_ERROR:   "✗",
    SEVERITY_WARNING: "⚠",
    SEVERITY_INFO:    "ℹ",
}
SEVERITY_ORDER = {SEVERITY_ERROR: 0, SEVERITY_WARNING: 1, SEVERITY_INFO: 2}


def print_report(report: Dict[str, Any], min_severity: str = "info") -> None:
    """Print a human-readable QA report to stdout."""
    if "error" in report:
        print(f"ERROR: {report['error']}", file=sys.stderr)
        return

    summary = report["summary"]
    print(f"\n{'='*60}")
    print(f"  PPTX QA Report: {report['file']}")
    print(f"{'='*60}")
    print(f"  Slides: {report['slides']}")
    print(f"  Issues: {summary['total']}  "
          f"[{summary['errors']} errors  "
          f"{summary['warnings']} warnings  "
          f"{summary['info']} info]")

    if summary["passed"]:
        print(f"\n  ✓ All checks passed!\n")
        return

    min_order = SEVERITY_ORDER.get(min_severity, 2)
    visible = [i for i in report["issues"] if SEVERITY_ORDER[i["severity"]] <= min_order]

    if not visible:
        print(f"\n  ✓ No issues at or above '{min_severity}' severity.\n")
        return

    print()
    # Group by slide
    by_slide: Dict[int, List] = {}
    for issue in visible:
        by_slide.setdefault(issue["slide"], []).append(issue)

    for slide_num in sorted(by_slide.keys()):
        print(f"  Slide {slide_num}:")
        for issue in by_slide[slide_num]:
            icon = SEVERITY_ICON[issue["severity"]]
            print(f"    {icon} [{issue['check']}] {issue['message']}")
            if "detail" in issue:
                print(f"        {issue['detail']}")
        print()

    print(f"{'='*60}\n")


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Automated QA checker for PPTX files.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Available checks:
  overflow     Text likely overflows its bounding box
  contrast     Low text/background contrast (WCAG AA)
  empty        Slides with no text or images
  placeholder  Leftover template placeholder text
  bullets      More than 6 bullet points on a slide
  fontsize     Font too small (body <12pt, title <20pt)
  offslide     Shapes positioned outside slide boundaries
  duplicates   Consecutive slides with identical titles
  titles       Content slides missing a title
  monotony     3+ consecutive slides with the same layout

Examples:
  python scripts/qa_check.py deck.pptx
  python scripts/qa_check.py deck.pptx --only overflow,placeholder
  python scripts/qa_check.py deck.pptx --min-severity warning
  python scripts/qa_check.py deck.pptx --output qa_report.json
  python scripts/qa_check.py deck.pptx --strict
        """
    )
    parser.add_argument("input", help="Input .pptx file to check")
    parser.add_argument("--output", "-o", help="Save JSON report to this file")
    parser.add_argument(
        "--only",
        help="Comma-separated list of checks to run (default: all)",
        default=None
    )
    parser.add_argument(
        "--min-severity",
        choices=["error", "warning", "info"],
        default="info",
        help="Minimum severity level to display (default: info)"
    )
    parser.add_argument(
        "--strict",
        action="store_true",
        help="Exit with code 1 if any errors or warnings are found"
    )
    parser.add_argument(
        "--json",
        action="store_true",
        help="Output raw JSON report instead of human-readable format"
    )

    args = parser.parse_args()

    enabled_checks = None
    if args.only:
        requested = [c.strip() for c in args.only.split(",")]
        invalid = [c for c in requested if c not in ALL_CHECKS]
        if invalid:
            print(f"Error: Unknown checks: {', '.join(invalid)}", file=sys.stderr)
            print(f"Available: {', '.join(ALL_CHECKS)}", file=sys.stderr)
            sys.exit(1)
        enabled_checks = requested

    report = run_qa(args.input, enabled_checks)

    if args.json:
        print(json.dumps(report, indent=2, ensure_ascii=False))
    else:
        print_report(report, min_severity=args.min_severity)

    if args.output:
        out_path = Path(args.output)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        with open(out_path, "w", encoding="utf-8") as f:
            json.dump(report, f, indent=2, ensure_ascii=False)
        print(f"JSON report saved to: {args.output}")

    if args.strict and report.get("summary", {}).get("errors", 0) > 0:
        sys.exit(1)
    if args.strict and report.get("summary", {}).get("warnings", 0) > 0:
        sys.exit(1)


if __name__ == "__main__":
    main()
